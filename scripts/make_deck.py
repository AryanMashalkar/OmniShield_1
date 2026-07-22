"""Generate the OmniShield pitch deck (OmniShield_Pitch.pptx) — neon SOC theme.

Uses assets/*.png (run scripts/gen_assets.py first) + baked-in slide transitions.
Run:  venv\\Scripts\\python scripts\\make_deck.py
"""
import os
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

A = os.path.join(os.path.dirname(__file__), "..", "assets")

BG = RGBColor(0x0A, 0x0F, 0x1E)
CARD = RGBColor(0x0D, 0x14, 0x24)
WHITE = RGBColor(0xF1, 0xF5, 0xF9)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
EMER = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def fullbg(s, name):
    s.shapes.add_picture(os.path.join(A, name), 0, 0, W, H)


def pic(s, name, l, t, w):
    return s.shapes.add_picture(os.path.join(A, name), Inches(l), Inches(t), width=Inches(w))


def box(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tf = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def para(tf, text, size, color, bold=False, first=False, align=PP_ALIGN.LEFT,
         space_after=8, bullet=None, bcolor=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    if bullet:
        r = p.add_run(); r.text = bullet + "  "
        r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = bcolor or BLUE
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = "Segoe UI"; r.font.color.rgb = color
    return p


def rect(s, l, t, w, h, fill=None, line=None, lw=1.5, rounded=False, shadow=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                             Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def chrome(s, kicker, kcolor, n):
    rect(s, 0.7, 0.72, 1.05, 0.13, fill=kcolor)
    para(box(s, 0.7, 0.9, 12, 0.5), kicker, 13, kcolor, bold=True, first=True)
    para(box(s, 0.7, 6.98, 9, 0.4),
         "OmniShield  ·  PS7: AI-Driven Cyber Resilience for Critical National Infrastructure",
         10, MUTED, first=True)
    para(box(s, 12.2, 6.98, 0.9, 0.4), str(n), 10, MUTED, first=True, align=PP_ALIGN.RIGHT)


def title(s, text, size=34, color=WHITE, top=1.3):
    para(box(s, 0.7, top, 12, 1.2), text, size, color, bold=True, first=True)


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def transition(s, kind="fade", spd="med", direction=None):
    sld = s._element
    for el in sld.findall(qn("p:transition")):
        sld.remove(el)
    tr = etree.SubElement(sld, qn("p:transition"))
    tr.set("spd", spd)
    child = etree.SubElement(tr, qn("p:" + kind))
    if direction:
        child.set("dir", direction)


# === 1. TITLE ==============================================================
s = slide(); fullbg(s, "hero_bg.png")
rect(s, 0.7, 2.25, 0.16, 2.75, fill=BLUE)
tf = box(s, 1.05, 2.2, 11.5, 3)
para(tf, "OmniShield", 66, WHITE, bold=True, first=True, space_after=4)
para(tf, "Autonomous, on-prem cyber resilience for", 22, MUTED, space_after=0)
para(tf, "Critical National Infrastructure", 22, BLUE, bold=True, space_after=16)
para(tf, "Predict.  Correlate.  Contain — before the breach.", 18, WHITE)
para(box(s, 1.05, 6.4, 11, 0.6),
     "Problem Statement 7   ·   100% local LLM   ·   Zero data egress", 13, EMER, first=True)
transition(s, "fade")
notes(s, "Open confident. Team name + PS7 on screen.")

# === 2. PROBLEM ============================================================
s = slide(); chrome(s, "THE PROBLEM", RED, 2)
title(s, "SOCs are reactive. Attackers are not.")
cards = [("1.59M", "CERT-In incidents, 2023", RED),
         ("2 weeks", "AIIMS Delhi ransomware outage", AMBER),
         (">70%", "govt entities on end-of-life IT", AMBER),
         ("weeks", "industry mean-time-to-detect", RED)]
x = 0.7
for big, sub, col in cards:
    rect(s, x, 2.6, 2.95, 1.9, fill=CARD, line=col, rounded=True)
    tf = box(s, x + 0.15, 2.8, 2.65, 1.6, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, big, 30, col, bold=True, first=True, align=PP_ALIGN.CENTER, space_after=6)
    para(tf, sub, 12.5, WHITE, align=PP_ALIGN.CENTER)
    x += 3.08
para(box(s, 0.7, 5.0, 12, 1.2),
     "Static rules + signatures are blind to zero-days and low-and-slow APTs. The data is there — the intelligence to act on it isn't.",
     18, MUTED, first=True)
transition(s, "push", direction="l")

# === 3. GAP ================================================================
s = slide(); chrome(s, "THE GAP", AMBER, 3)
title(s, "The data was always there.\nThe intelligence layer wasn't.", size=32)
tf = box(s, 0.7, 3.7, 12, 2.4)
para(tf, "We rebuilt the SOC to be predictive and agentic — not reactive and rule-bound.",
     23, WHITE, first=True, space_after=16)
para(tf, "Detect from behaviour  →  attribute with grounded AI  →  forecast the next move  →  contain in one click.",
     17, BLUE)
transition(s, "fade")

# === 4. CLOSED LOOP ========================================================
s = slide(); chrome(s, "WHAT OMNISHIELD IS", BLUE, 4)
title(s, "One closed, autonomous loop")
pic(s, "loop.png", 8.0, 1.7, 5.0)
tf = box(s, 0.7, 2.7, 7.0, 4)
para(tf, "Six cooperating agents, one loop:", 18, WHITE, bold=True, first=True, space_after=14)
for t, c in [("DETECT — behavioural UEBA", BLUE), ("ATTRIBUTE — MITRE ATT&CK RAG", PURPLE),
             ("PREDICT — next-move forecast", AMBER), ("CORRELATE — group by host", EMER),
             ("CONTAIN — human-gated SOAR", RED), ("AUDIT — SQLite log", MUTED)]:
    para(tf, t, 16, WHITE, bullet="▸", bcolor=c, space_after=10)
para(box(s, 0.7, 6.1, 7, 0.8), "100% local. Air-gap capable. Nothing leaves the perimeter.",
     15, EMER, first=True)
transition(s, "push", direction="l")

# === 5. DATA ===============================================================
s = slide(); chrome(s, "DIFFERENTIATOR · MODERN DATA", EMER, 5)
title(s, "Not 1998 data. Real modern enterprise traffic.")
tf = box(s, 0.7, 2.5, 12, 4)
para(tf, "Live WebSocket stream of CIC-IDS-2017 — authentic 2017 enterprise flows.",
     20, WHITE, first=True, bullet="▸", bcolor=EMER, space_after=14)
para(tf, "Modern threats: Heartbleed · Botnet · DDoS · PortScan · Slowloris · Infiltration.",
     20, WHITE, bullet="▸", bcolor=EMER, space_after=14)
para(tf, "Dataset-pluggable: NSL-KDD · UNSW-NB15 · CIC-IDS-2017 — switched by one env var.",
     20, WHITE, bullet="▸", bcolor=EMER, space_after=14)
para(tf, "Most teams benchmark on 1998 KDD. We stream the traffic real SOCs actually face.",
     17, AMBER)
transition(s, "fade")
notes(s, "ACCURACY: say 'authentic CIC-IDS-2017' only if the real CSV is loaded (scripts/prepare_cic.py). "
        "Otherwise say 'CIC-IDS-2017-format modern flows'.")

# === 6. ML ENGINE ==========================================================
s = slide(); chrome(s, "ML ENGINE · BEHAVIOURAL DETECTION (UEBA)", BLUE, 6)
title(s, "Multivariate IsolationForest — no signatures")
tf = box(s, 0.7, 2.35, 6.4, 3)
para(tf, "One-class IsolationForest trains on baseline normal flows.",
     17, WHITE, first=True, bullet="▸", bcolor=BLUE, space_after=12)
para(tf, "Per-entity UEBA: every host scored against ITS OWN normal — not a global threshold.",
     17, WHITE, bullet="▸", bcolor=BLUE, space_after=12)
para(tf, "Catches the low-and-slow scans the univariate baseline can't see.",
     17, WHITE, bullet="▸", bcolor=BLUE, space_after=12)
pic(s, "gauge.png", 0.9, 4.9, 3.4)
rect(s, 7.35, 2.2, 5.6, 4.35, fill=CARD, line=EMER, rounded=True)
pic(s, "detection_chart.png", 7.55, 2.45, 5.2)
transition(s, "fade")
notes(s, "F1 0.93 = NSL-KDD one-class benchmark (python -m evaluation.eval_detection). Reproducible.")

# === 7. ATTRIBUTION ========================================================
s = slide(); chrome(s, "AI ATTRIBUTION · GROUNDED & SOVEREIGN", PURPLE, 7)
title(s, "Local LLM, grounded in MITRE ATT&CK")
tf = box(s, 0.7, 2.5, 6.4, 4)
para(tf, "Local Llama 3.1 (8B) + RAG over 222 MITRE ATT&CK techniques.",
     18, WHITE, first=True, bullet="▸", bcolor=PURPLE, space_after=14)
para(tf, "Retrieval grounding doubles accuracy — an honest ablation.",
     18, WHITE, bullet="▸", bcolor=PURPLE, space_after=14)
para(tf, "~7s, on-box. No cloud, no API keys — nothing leaves the perimeter.",
     18, WHITE, bullet="▸", bcolor=PURPLE, space_after=14)
rect(s, 7.35, 2.2, 5.6, 4.35, fill=CARD, line=PURPLE, rounded=True)
pic(s, "attribution_chart.png", 7.55, 2.45, 5.2)
transition(s, "fade")

# === 8. NEXT-MOVE ==========================================================
s = slide(); chrome(s, "THE INNOVATION · NEXT-MOVE PREDICTOR", AMBER, 8)
title(s, "We forecast the attacker's next move")
para(box(s, 0.7, 2.25, 12, 1.0),
     "A deterministic graph over the ATT&CK kill chain maps behaviour to a tactic — then predicts the next stage, so you pre-stage containment before it happens.",
     16, MUTED, first=True)
pic(s, "killchain.png", 0.5, 3.3, 12.4)
para(box(s, 0.7, 6.2, 12, 0.7),
     "PortScan → Discovery → forecast: Lateral Movement.   DDoS → Impact → \"terminal — isolate now.\"",
     15, AMBER, first=True, bold=True)
transition(s, "push", direction="l")
notes(s, "Your wow moment — slow down. 'We don't just say what's happening; we say what's next.'")

# === 9. CORRELATION ========================================================
s = slide(); chrome(s, "ACTIVE INCIDENT CORRELATION", EMER, 9)
title(s, "Weak signals → one story. No alert fatigue.")
para(box(s, 0.7, 2.4, 12, 1.1),
     "Scattered anomalies are grouped by source host into a single multi-stage incident — 40 red rows become one narrative.",
     18, WHITE, first=True)
rect(s, 0.8, 3.9, 11.7, 1.9, fill=CARD, line=AMBER, rounded=True)
tf = box(s, 1.1, 4.05, 11.1, 1.6, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "host  192.168.1.66", 18, BLUE, bold=True, first=True, space_after=10)
para(tf, "Discovery  →  Credential Access  →  Lateral Movement  →  Exfiltration      ⟶   forecast: IMPACT",
     18, WHITE)
transition(s, "fade")

# === 10. SOAR ==============================================================
s = slide(); chrome(s, "SAFE AUTONOMY · HUMAN-GATED SOAR", RED, 10)
title(s, "One click. Threat neutralized. Fully audited.")
tf = box(s, 0.7, 2.6, 12, 4)
para(tf, "Analyst clicks once → firewall rule injected → source IP isolated → NEUTRALIZED.",
     19, WHITE, first=True, bullet="▸", bcolor=RED, space_after=15)
para(tf, "The AI never blocks alone: HTTP 428 until a human authorizes (blast-radius gate).",
     19, WHITE, bullet="▸", bcolor=RED, space_after=15)
para(tf, "Every action logged to SQLite — reversible and court-auditable.",
     19, WHITE, bullet="▸", bcolor=RED, space_after=15)
para(tf, "The AI recommends.  A human authorizes.  The system acts and audits.", 19, EMER)
transition(s, "fade")

# === 11. ARCHITECTURE ======================================================
s = slide(); chrome(s, "ARCHITECTURE & PROOF", BLUE, 11)
title(s, "Engineered, not hand-waved")
tf = box(s, 0.7, 2.6, 12, 4)
for t in ["FastAPI + WebSocket · React dashboard · Chroma vector store · Ollama local LLM.",
          "Modular agents: Detection · Attribution · Forecast · Response · Audit.",
          "47 automated tests (pytest). Reproducible benchmarks in evaluation/results/.",
          "Hardened: no hard-coded secrets, CORS allow-list, prompt-injection tested."]:
    para(tf, t, 18, WHITE, first=t.startswith("FastAPI"), bullet="▸", bcolor=BLUE, space_after=16)
transition(s, "fade")

# === 12. IMPACT ============================================================
s = slide(); chrome(s, "BUSINESS IMPACT", EMER, 12)
title(s, "Weeks → seconds, on your own hardware")
cards = [("MTTD / MTTR", "weeks → <60s", "on confirmed anomalies", EMER),
         ("DATA SOVEREIGNTY", "zero egress", "on-prem, air-gap capable", BLUE),
         ("ADDRESSABLE", "govt CNI", "EOL-IT estate, no cloud SIEM", AMBER)]
x = 0.7
for head, big, sub, col in cards:
    rect(s, x, 2.9, 3.95, 2.4, fill=CARD, line=col, rounded=True)
    tf = box(s, x + 0.2, 3.1, 3.55, 2.0, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, head, 13, MUTED, bold=True, first=True, align=PP_ALIGN.CENTER, space_after=8)
    para(tf, big, 25, col, bold=True, align=PP_ALIGN.CENTER, space_after=6)
    para(tf, sub, 13, WHITE, align=PP_ALIGN.CENTER)
    x += 4.1
para(box(s, 0.7, 5.7, 12, 1),
     "IBM: avg breach ~$4.4M; fast containment saves ~$1M+. AIIMS lost 2 weeks. We isolate in under a minute.",
     15, MUTED, first=True)
transition(s, "push", direction="l")

# === 13. CLOSE =============================================================
s = slide(); fullbg(s, "hero_bg.png")
rect(s, 0.7, 2.1, 0.16, 2.9, fill=EMER)
tf = box(s, 1.05, 2.0, 11.5, 3.4)
para(tf, "OmniShield", 50, WHITE, bold=True, first=True, space_after=8)
para(tf, "Detect from behaviour.  Predict the next move.  Contain in one click.  On-prem.",
     20, EMER, space_after=18)
para(tf, "It doesn't just watch the breach — it gets ahead of it. And it runs live.", 18, WHITE)
para(box(s, 1.05, 6.3, 11, 0.8),
     "github.com/<your-repo>   ·   Team <name>   ·   Problem Statement 7", 14, BLUE, first=True)
transition(s, "fade")
notes(s, "End on 'and it runs live — let me show you' → cut to the demo.")

out = "OmniShield_Pitch.pptx"
try:
    prs.save(out)
except PermissionError:
    out = "OmniShield_Pitch_neon.pptx"
    prs.save(out)
    print("[warn] OmniShield_Pitch.pptx was locked (close PowerPoint); wrote", out, "instead.")
# Re-open to validate the XML (transitions) is well-formed.
Presentation(out)
print(f"Saved + validated {out} · {len(prs.slides._sldIdLst)} slides")
